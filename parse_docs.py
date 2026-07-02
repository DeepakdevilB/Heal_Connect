import os
import sys

def parse_docx(file_path):
    import docx
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_xlsx(file_path):
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text.append(f"Sheet: {sheet}")
        for row in ws.iter_rows(values_only=True):
            text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
    return "\n".join(text)

def parse_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

if __name__ == "__main__":
    folder = "d:\\Tara_InfoTech_Projects\\Heal_Connect"
    with open(os.path.join(folder, "extracted_texts.txt"), "w", encoding="utf-8") as out:
        for f in os.listdir(folder):
            if f.endswith(".docx") and not f.startswith("~"):
                out.write(f"=== {f} ===\n")
                try:
                    out.write(parse_docx(os.path.join(folder, f)) + "\n")
                except Exception as e:
                    out.write(f"Error reading docx: {e}\n")
            elif f.endswith(".xlsx") and not f.startswith("~"):
                out.write(f"=== {f} ===\n")
                try:
                    out.write(parse_xlsx(os.path.join(folder, f)) + "\n")
                except Exception as e:
                    out.write(f"Error reading xlsx: {e}\n")
            elif f.endswith(".pptx") and not f.startswith("~"):
                out.write(f"=== {f} ===\n")
                try:
                    out.write(parse_pptx(os.path.join(folder, f)) + "\n")
                except Exception as e:
                    out.write(f"Error reading pptx: {e}\n")
