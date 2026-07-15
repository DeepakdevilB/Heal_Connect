import os

def parse_docx(file_path):
    import docx
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_xlsx(file_path):
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text.append(f"Sheet: {sheet}")
        for row in ws.iter_rows(values_only=True):
            text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
    return "\n".join(text)

def parse_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

if __name__ == "__main__":
    folder = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # project root
    output_path = os.path.join(folder, "extracted_texts.txt")

    with open(output_path, "w", encoding="utf-8") as out:
        for f in sorted(os.listdir(folder)):
            full_path = os.path.join(folder, f)
            if f.startswith("~") or not os.path.isfile(full_path):
                continue
            if f.endswith(".docx"):
                out.write(f"=== {f} ===\n")
                try:
                    out.write(parse_docx(full_path) + "\n")
                except Exception as e:
                    out.write(f"Error reading docx: {e}\n")
            elif f.endswith(".xlsx"):
                out.write(f"=== {f} ===\n")
                try:
                    out.write(parse_xlsx(full_path) + "\n")
                except Exception as e:
                    out.write(f"Error reading xlsx: {e}\n")
            elif f.endswith(".pptx"):
                out.write(f"=== {f} ===\n")
                try:
                    out.write(parse_pptx(full_path) + "\n")
                except Exception as e:
                    out.write(f"Error reading pptx: {e}\n")

    print(f"Done! Output saved to: {output_path}")
